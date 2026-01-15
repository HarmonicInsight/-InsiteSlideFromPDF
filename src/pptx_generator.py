"""
PowerPoint Generation Module
Uses python-pptx to create PowerPoint presentations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image
import io
from typing import List, Dict, Tuple, Optional, Union
from dataclasses import dataclass
from enum import Enum


class ShapeType(Enum):
    """Supported shape types."""
    RECTANGLE = "rectangle"
    ROUNDED_RECTANGLE = "rounded_rectangle"
    OVAL = "oval"
    CIRCLE = "circle"
    TRIANGLE = "triangle"
    ARROW_RIGHT = "arrow_right"
    ARROW_LEFT = "arrow_left"
    ARROW_UP = "arrow_up"
    ARROW_DOWN = "arrow_down"
    STAR = "star"
    DIAMOND = "diamond"
    PENTAGON = "pentagon"
    HEXAGON = "hexagon"


# Mapping from ShapeType to MSO_SHAPE
SHAPE_TYPE_MAP = {
    ShapeType.RECTANGLE: MSO_SHAPE.RECTANGLE,
    ShapeType.ROUNDED_RECTANGLE: MSO_SHAPE.ROUNDED_RECTANGLE,
    ShapeType.OVAL: MSO_SHAPE.OVAL,
    ShapeType.CIRCLE: MSO_SHAPE.OVAL,
    ShapeType.TRIANGLE: MSO_SHAPE.ISOSCELES_TRIANGLE,
    ShapeType.ARROW_RIGHT: MSO_SHAPE.RIGHT_ARROW,
    ShapeType.ARROW_LEFT: MSO_SHAPE.LEFT_ARROW,
    ShapeType.ARROW_UP: MSO_SHAPE.UP_ARROW,
    ShapeType.ARROW_DOWN: MSO_SHAPE.DOWN_ARROW,
    ShapeType.STAR: MSO_SHAPE.STAR_5_POINT,
    ShapeType.DIAMOND: MSO_SHAPE.DIAMOND,
    ShapeType.PENTAGON: MSO_SHAPE.PENTAGON,
    ShapeType.HEXAGON: MSO_SHAPE.HEXAGON,
}


@dataclass
class SlideElement:
    """Base class for slide elements."""
    x: float  # Position in inches
    y: float
    width: float
    height: float


@dataclass
class TextElement(SlideElement):
    """Text element for slides."""
    text: str
    font_size: float = 12
    font_name: str = "Arial"
    font_color: Tuple[int, int, int] = (0, 0, 0)
    bold: bool = False
    italic: bool = False
    alignment: str = "left"  # left, center, right


@dataclass
class ShapeElement(SlideElement):
    """Shape element for slides."""
    shape_type: ShapeType = ShapeType.RECTANGLE
    fill_color: Optional[Tuple[int, int, int]] = None
    line_color: Tuple[int, int, int] = (0, 0, 0)
    line_width: float = 1.0
    text: Optional[str] = None
    text_color: Tuple[int, int, int] = (0, 0, 0)
    text_size: float = 10


@dataclass
class ImageElement(SlideElement):
    """Image element for slides."""
    image: Union[Image.Image, str]  # PIL Image or file path


@dataclass
class SlideContent:
    """Contains all elements for a single slide."""
    title: Optional[str] = None
    elements: List[Union[TextElement, ShapeElement, ImageElement]] = None

    def __post_init__(self):
        if self.elements is None:
            self.elements = []


class PPTXGenerator:
    """Generates PowerPoint presentations."""

    # Standard slide dimensions (16:9 aspect ratio)
    SLIDE_WIDTH_INCHES = 13.333
    SLIDE_HEIGHT_INCHES = 7.5

    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize PowerPoint generator.

        Args:
            template_path: Optional path to a template PPTX file.
        """
        if template_path:
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            # Set slide dimensions
            self.prs.slide_width = Inches(self.SLIDE_WIDTH_INCHES)
            self.prs.slide_height = Inches(self.SLIDE_HEIGHT_INCHES)

    def add_blank_slide(self) -> int:
        """
        Add a blank slide.

        Returns:
            Index of the new slide.
        """
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)
        return len(self.prs.slides) - 1

    def add_title_slide(self, title: str, subtitle: Optional[str] = None) -> int:
        """
        Add a title slide.

        Args:
            title: Slide title.
            subtitle: Optional subtitle.

        Returns:
            Index of the new slide.
        """
        title_layout = self.prs.slide_layouts[0]  # Title layout
        slide = self.prs.slides.add_slide(title_layout)

        if slide.shapes.title:
            slide.shapes.title.text = title

        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

        return len(self.prs.slides) - 1

    def add_slide_with_content(self, content: SlideContent) -> int:
        """
        Add a slide with specified content.

        Args:
            content: SlideContent object with elements to add.

        Returns:
            Index of the new slide.
        """
        slide_idx = self.add_blank_slide()
        slide = self.prs.slides[slide_idx]

        # Add title if specified
        if content.title:
            self.add_text_box(
                slide_idx,
                TextElement(
                    x=0.5, y=0.3,
                    width=12.33, height=0.8,
                    text=content.title,
                    font_size=28,
                    bold=True,
                    alignment="center"
                )
            )

        # Add elements
        for element in content.elements:
            if isinstance(element, TextElement):
                self.add_text_box(slide_idx, element)
            elif isinstance(element, ShapeElement):
                self.add_shape(slide_idx, element)
            elif isinstance(element, ImageElement):
                self.add_image(slide_idx, element)

        return slide_idx

    def add_text_box(self, slide_idx: int, element: TextElement):
        """
        Add a text box to a slide.

        Args:
            slide_idx: Slide index.
            element: TextElement to add.
        """
        slide = self.prs.slides[slide_idx]

        left = Inches(element.x)
        top = Inches(element.y)
        width = Inches(element.width)
        height = Inches(element.height)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        frame = textbox.text_frame
        frame.word_wrap = True

        p = frame.paragraphs[0]
        p.text = element.text

        # Set alignment
        if element.alignment == "center":
            p.alignment = PP_ALIGN.CENTER
        elif element.alignment == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        # Set font properties
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = element.font_name
        run.font.size = Pt(element.font_size)
        run.font.bold = element.bold
        run.font.italic = element.italic
        run.font.color.rgb = RGBColor(*element.font_color)

    def add_shape(self, slide_idx: int, element: ShapeElement):
        """
        Add a shape to a slide.

        Args:
            slide_idx: Slide index.
            element: ShapeElement to add.
        """
        slide = self.prs.slides[slide_idx]

        left = Inches(element.x)
        top = Inches(element.y)
        width = Inches(element.width)
        height = Inches(element.height)

        # Get MSO_SHAPE type
        mso_shape = SHAPE_TYPE_MAP.get(element.shape_type, MSO_SHAPE.RECTANGLE)

        shape = slide.shapes.add_shape(mso_shape, left, top, width, height)

        # Set fill color
        if element.fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*element.fill_color)
        else:
            shape.fill.background()

        # Set line properties
        shape.line.color.rgb = RGBColor(*element.line_color)
        shape.line.width = Pt(element.line_width)

        # Add text if specified
        if element.text:
            shape.text = element.text
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                if p.runs:
                    p.runs[0].font.size = Pt(element.text_size)
                    p.runs[0].font.color.rgb = RGBColor(*element.text_color)

    def add_image(self, slide_idx: int, element: ImageElement):
        """
        Add an image to a slide.

        Args:
            slide_idx: Slide index.
            element: ImageElement to add.
        """
        slide = self.prs.slides[slide_idx]

        left = Inches(element.x)
        top = Inches(element.y)
        width = Inches(element.width)
        height = Inches(element.height)

        if isinstance(element.image, str):
            # File path
            slide.shapes.add_picture(element.image, left, top, width, height)
        else:
            # PIL Image
            img_bytes = io.BytesIO()
            element.image.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            slide.shapes.add_picture(img_bytes, left, top, width, height)

    def add_image_full_slide(self, slide_idx: int, image: Union[Image.Image, str]):
        """
        Add an image that fills the entire slide.

        Args:
            slide_idx: Slide index.
            image: PIL Image or file path.
        """
        element = ImageElement(
            x=0, y=0,
            width=self.SLIDE_WIDTH_INCHES,
            height=self.SLIDE_HEIGHT_INCHES,
            image=image
        )
        self.add_image(slide_idx, element)

    def save(self, output_path: str):
        """
        Save the presentation to a file.

        Args:
            output_path: Path to save the PPTX file.
        """
        self.prs.save(output_path)

    def get_slide_count(self) -> int:
        """Get the number of slides."""
        return len(self.prs.slides)

    @staticmethod
    def convert_position_to_inches(
        x: float, y: float,
        source_width: float, source_height: float,
        target_width: float = 13.333,
        target_height: float = 7.5
    ) -> Tuple[float, float]:
        """
        Convert position from source dimensions to inches.

        Args:
            x, y: Position in source units.
            source_width, source_height: Source dimensions.
            target_width, target_height: Target dimensions in inches.

        Returns:
            Position in inches.
        """
        x_inches = (x / source_width) * target_width
        y_inches = (y / source_height) * target_height
        return x_inches, y_inches

    @staticmethod
    def convert_size_to_inches(
        width: float, height: float,
        source_width: float, source_height: float,
        target_width: float = 13.333,
        target_height: float = 7.5
    ) -> Tuple[float, float]:
        """
        Convert size from source dimensions to inches.

        Args:
            width, height: Size in source units.
            source_width, source_height: Source dimensions.
            target_width, target_height: Target dimensions in inches.

        Returns:
            Size in inches.
        """
        w_inches = (width / source_width) * target_width
        h_inches = (height / source_height) * target_height
        return w_inches, h_inches


def create_presentation_from_images(
    images: List[Union[Image.Image, str]],
    output_path: str,
    progress_callback=None
):
    """
    Create a presentation from a list of images (one per slide).

    Args:
        images: List of PIL Images or file paths.
        output_path: Path to save the PPTX file.
        progress_callback: Optional callback(current, total).
    """
    generator = PPTXGenerator()
    total = len(images)

    for i, image in enumerate(images):
        slide_idx = generator.add_blank_slide()
        generator.add_image_full_slide(slide_idx, image)

        if progress_callback:
            progress_callback(i + 1, total)

    generator.save(output_path)


def create_slide_from_content(
    generator: PPTXGenerator,
    title: Optional[str],
    texts: List[Dict],
    shapes: List[Dict],
    images: List[Dict],
    source_width: float,
    source_height: float
) -> int:
    """
    Create a slide from extracted content.

    Args:
        generator: PPTXGenerator instance.
        title: Optional slide title.
        texts: List of text dictionaries with bbox and text.
        shapes: List of shape dictionaries with bbox and color.
        images: List of image dictionaries with bbox and image.
        source_width, source_height: Source dimensions for coordinate conversion.

    Returns:
        Slide index.
    """
    elements = []

    # Add text elements
    for text_data in texts:
        x, y, w, h = text_data.get('bbox', (0, 0, 100, 20))
        x_in, y_in = PPTXGenerator.convert_position_to_inches(
            x, y, source_width, source_height
        )
        w_in, h_in = PPTXGenerator.convert_size_to_inches(
            w, h, source_width, source_height
        )

        elements.append(TextElement(
            x=x_in, y=y_in,
            width=max(w_in, 0.5), height=max(h_in, 0.3),
            text=text_data.get('text', ''),
            font_size=text_data.get('font_size', 12),
            font_name=text_data.get('font_name', 'Arial'),
            font_color=text_data.get('color', (0, 0, 0))
        ))

    # Add shape elements
    for shape_data in shapes:
        x, y, w, h = shape_data.get('bbox', (0, 0, 50, 50))
        x_in, y_in = PPTXGenerator.convert_position_to_inches(
            x, y, source_width, source_height
        )
        w_in, h_in = PPTXGenerator.convert_size_to_inches(
            w, h, source_width, source_height
        )

        shape_type_str = shape_data.get('shape_type', 'rectangle')
        shape_type = ShapeType.RECTANGLE
        for st in ShapeType:
            if st.value == shape_type_str:
                shape_type = st
                break

        elements.append(ShapeElement(
            x=x_in, y=y_in,
            width=max(w_in, 0.2), height=max(h_in, 0.2),
            shape_type=shape_type,
            fill_color=shape_data.get('fill_color'),
            line_color=shape_data.get('line_color', (0, 0, 0)),
            text=shape_data.get('text')
        ))

    # Add image elements
    for img_data in images:
        x, y, w, h = img_data.get('bbox', (0, 0, 100, 100))
        x_in, y_in = PPTXGenerator.convert_position_to_inches(
            x, y, source_width, source_height
        )
        w_in, h_in = PPTXGenerator.convert_size_to_inches(
            w, h, source_width, source_height
        )

        elements.append(ImageElement(
            x=x_in, y=y_in,
            width=max(w_in, 0.5), height=max(h_in, 0.5),
            image=img_data.get('image')
        ))

    content = SlideContent(title=title, elements=elements)
    return generator.add_slide_with_content(content)
